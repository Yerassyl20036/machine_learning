#!/usr/bin/env python3
"""
Knee MRI Segmentation — PowerPoint Presentation (in Russian).
Usage: python generate_slides.py
Output: knee_mri_presentation_easy_explained.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── Paths ───────────────────────────────────────────────────────────────────
FIGURES = Path(__file__).parent / "results" / "eda_figures"

# ─── Color palette ───────────────────────────────────────────────────────────
C_BG    = RGBColor(0x12, 0x22, 0x38)
C_CARD  = RGBColor(0x1C, 0x33, 0x4D)
C_BLUE  = RGBColor(0x00, 0x99, 0xFF)
C_GREEN = RGBColor(0x00, 0xCC, 0x99)
C_ORANGE= RGBColor(0xFF, 0x99, 0x00)
C_RED   = RGBColor(0xFF, 0x44, 0x55)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LGRAY = RGBColor(0xCC, 0xDD, 0xEE)
C_DGRAY = RGBColor(0x77, 0x99, 0xBB)

# ─── Slide dimensions (16:9 widescreen) ──────────────────────────────────────
SW = Inches(13.33)
SH = Inches(7.5)

TOTAL = 13


# ═════════════════════════════════════════════════════════════════════════════
# Primitive helpers
# ═════════════════════════════════════════════════════════════════════════════

def bg(slide, color=C_BG):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def rect(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def txt(slide, content, l, t, w, h, size=18, bold=False,
        color=C_WHITE, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = content
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def bullets(slide, lines, l, t, w, h, size=16, color=C_LGRAY):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(7)
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = "●  " + line
        r.font.size = Pt(size)
        r.font.color.rgb = color


def img(slide, path, l, t, w, h):
    p = Path(path)
    if p.exists():
        slide.shapes.add_picture(str(p), l, t, w, h)
    else:
        rect(slide, l, t, w, h, C_CARD)
        txt(slide, f"[ {p.name} ]", l, t + h // 3, w, Inches(0.4),
            size=11, color=C_DGRAY, align=PP_ALIGN.CENTER)


def header(slide, title, sub=None):
    rect(slide, 0, 0, SW, Inches(1.05), C_CARD)
    rect(slide, 0, 0, Inches(0.1), Inches(1.05), C_BLUE)
    txt(slide, title, Inches(0.25), Inches(0.1), Inches(12), Inches(0.6),
        size=28, bold=True)
    if sub:
        txt(slide, sub, Inches(0.25), Inches(0.68), Inches(12), Inches(0.35),
            size=14, color=C_BLUE, italic=True)


def footer(slide, num):
    rect(slide, 0, SH - Inches(0.07), SW, Inches(0.07), C_BLUE)
    txt(slide, f"{num} / {TOTAL}",
        SW - Inches(1.0), SH - Inches(0.38), Inches(0.85), Inches(0.3),
        size=11, color=C_DGRAY, align=PP_ALIGN.RIGHT)


def card(slide, l, t, w, h, accent=None):
    rect(slide, l, t, w, h, C_CARD)
    if accent:
        rect(slide, l, t, Inches(0.07), h, accent)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE  1 — Title
# ═════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)

    # Left panel
    rect(s, 0, 0, Inches(7.0), SH, C_CARD)
    rect(s, 0, 0, Inches(0.12), SH, C_BLUE)

    txt(s, "Сегментация МРТ\nколенного сустава",
        Inches(0.3), Inches(0.8), Inches(6.4), Inches(2.2),
        size=36, bold=True)

    txt(s, "Автоматическое определение\nкостных и хрящевых структур\nпо снимкам МРТ",
        Inches(0.3), Inches(3.1), Inches(6.4), Inches(1.4),
        size=18, color=C_LGRAY)

    rect(s, Inches(0.3), Inches(4.65), Inches(5.5), Inches(0.05), C_GREEN)

    txt(s, "Маратов Ерасыл Балканович",
        Inches(0.3), Inches(4.8), Inches(6.0), Inches(0.45),
        size=18, bold=True, color=C_GREEN)

    txt(s, "7М06101  |  AIU  |  2025–2026  |  Алматы",
        Inches(0.3), Inches(5.35), Inches(6.0), Inches(0.4),
        size=13, color=C_DGRAY)

    # Right panel: MRI sample grid
    img(s, FIGURES / "sample_grid_kl_grades.png",
        Inches(7.2), Inches(0.2), Inches(5.9), Inches(7.0))

    footer(s, 1)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE  2 — Problem
# ═════════════════════════════════════════════════════════════════════════════

def slide_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Зачем это нужно?")
    footer(s, 2)

    # Top banner
    card(s, Inches(0.3), Inches(1.2), Inches(12.7), Inches(1.5), C_ORANGE)
    txt(s, "Остеоартрит — самое распространённое заболевание суставов. 520 млн человек в мире.",
        Inches(0.55), Inches(1.35), Inches(12.0), Inches(1.1),
        size=20, bold=True)

    # Problem card
    card(s, Inches(0.3), Inches(2.9), Inches(6.1), Inches(4.1), C_ORANGE)
    txt(s, "Проблема сегодня",
        Inches(0.55), Inches(2.98), Inches(5.7), Inches(0.45),
        size=18, bold=True, color=C_ORANGE)
    bullets(s, [
        "Врач вручную обводит ткани — 30–90 минут на одного пациента",
        "Разные врачи могут дать разные оценки одного снимка",
        "Чем позже находят болезнь — тем хуже лечение",
    ], Inches(0.55), Inches(3.5), Inches(5.7), Inches(3.2), size=15)

    # Solution card
    card(s, Inches(6.6), Inches(2.9), Inches(6.4), Inches(4.1), C_GREEN)
    txt(s, "Что мы делаем",
        Inches(6.85), Inches(2.98), Inches(6.0), Inches(0.45),
        size=18, bold=True, color=C_GREEN)
    bullets(s, [
        "Алгоритм сам находит кость и хрящ на МРТ-снимке",
        "Работает за секунды вместо часа",
        "Результат одинаковый при каждом запуске — нет субъективности",
        "Помогает врачу принять решение быстрее",
    ], Inches(6.85), Inches(3.5), Inches(6.0), Inches(3.2), size=15)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE  3 — Data + KL grade explanation
# ═════════════════════════════════════════════════════════════════════════════

def slide_data(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Данные и шкала Kellgren–Lawrence (KL)",
           "Что мы использовали и как читать диагноз")
    footer(s, 3)

    # Stats bar
    stats = [
        (C_BLUE,   "Источник",  "Kaggle\n(открытый датасет)"),
        (C_GREEN,  "Снимков",   "Тысячи 2D-срезов\nМРТ коленного сустава"),
        (C_ORANGE, "Классов",   "KL-0 ... KL-4\n5 стадий артрита"),
        (C_DGRAY,  "Разбивка",  "Train / Val / Test"),
    ]
    for i, (col, label, val) in enumerate(stats):
        lx = Inches(0.3) + i * Inches(3.2)
        rect(s, lx, Inches(1.2), Inches(3.0), Inches(1.6), C_CARD)
        rect(s, lx, Inches(1.2), Inches(3.0), Inches(0.1), col)
        txt(s, label, lx, Inches(1.35), Inches(3.0), Inches(0.4),
            size=13, color=col, align=PP_ALIGN.CENTER)
        txt(s, val, lx, Inches(1.72), Inches(3.0), Inches(1.0),
            size=15, bold=True, align=PP_ALIGN.CENTER)

    # KL legend
    kl_rows = [
        (C_GREEN,                    "KL-0", "Норма — сустав здоровый"),
        (C_BLUE,                     "KL-1", "Ранние изменения (едва заметны)"),
        (C_ORANGE,                   "KL-2", "Умеренная стадия — изменения заметны"),
        (RGBColor(0xFF, 0x66, 0x33), "KL-3", "Выраженная стадия"),
        (C_RED,                      "KL-4", "Тяжёлая стадия артрита"),
    ]
    card(s, Inches(0.3), Inches(2.95), Inches(5.5), Inches(4.1))
    txt(s, "Шкала KL — это не разные болезни, а стадии одной:",
        Inches(0.5), Inches(3.05), Inches(5.0), Inches(0.45),
        size=16, bold=True, color=C_LGRAY)
    y = Inches(3.6)
    for col, stage, desc in kl_rows:
        rect(s, Inches(0.5), y, Inches(0.65), Inches(0.45), col)
        txt(s, stage, Inches(0.5), y + Inches(0.04), Inches(0.65), Inches(0.38),
            size=12, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
        txt(s, desc, Inches(1.25), y + Inches(0.04), Inches(4.4), Inches(0.42),
            size=14, color=C_LGRAY)
        y += Inches(0.62)

    txt(s, "KL-0 = здоровый, KL-4 = самый тяжёлый. Модель выбирает наиболее вероятную стадию.",
        Inches(0.3), Inches(7.1), Inches(5.5), Inches(0.3),
        size=10, color=C_DGRAY, italic=True)

    # MRI sample image
    img(s, FIGURES / "sample_grid_kl_grades.png",
        Inches(6.0), Inches(2.95), Inches(7.1), Inches(4.1))
    txt(s, "Примеры МРТ: KL-0 (норма) → KL-4 (тяжёлый артрит)",
        Inches(6.0), Inches(7.1), Inches(7.1), Inches(0.3),
        size=11, color=C_DGRAY, italic=True, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE  4 — Откуда 20 параметров + почему Random Forest
# ═════════════════════════════════════════════════════════════════════════════

def slide_params_and_rf(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Откуда 20 параметров и почему Random Forest")
    footer(s, 4)

    # Left: where do 20 parameters come from
    card(s, Inches(0.3), Inches(1.2), Inches(6.2), Inches(5.7), C_BLUE)
    txt(s, "Откуда 20 параметров",
        Inches(0.55), Inches(1.35), Inches(5.7), Inches(0.45),
        size=19, bold=True, color=C_BLUE)
    bullets(s, [
        "Берём МРТ-снимок как обычную картинку",
        "Считаем числа из пикселей: яркость, резкость, текстура",
        "Добавляем геометрию: площадь кости, ширина щели, хрящ",
        "Итого: 20 чисел на каждый снимок → строка в таблице",
        "Это не анкета, не ввод вручную — только математика из снимка",
    ], Inches(0.55), Inches(1.95), Inches(5.7), Inches(4.7), size=15)

    # Right: why Random Forest
    card(s, Inches(6.8), Inches(1.2), Inches(6.2), Inches(5.7), C_GREEN)
    txt(s, "Почему Random Forest",
        Inches(7.05), Inches(1.35), Inches(5.7), Inches(0.45),
        size=19, bold=True, color=C_GREEN)
    bullets(s, [
        "Данные нелинейные (это видно на PCA-графике)",
        "Хорошо работает на табличных данных без сложной настройки",
        "Показывает важность признаков — видно, что влияет на диагноз",
        "Устойчив к выбросам и дисбалансу классов",
        "Стабильнее простой линейной модели (проверено сравнением)",
    ], Inches(7.05), Inches(1.95), Inches(5.7), Inches(4.7), size=15)

    # Bottom summary bar
    card(s, Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.3), C_ORANGE)
    txt(s, "Итого: снимок  →  20 параметров из пикселей  →  Random Forest  →  KL-стадия",
        Inches(0.55), Inches(7.12), Inches(12.2), Inches(0.24),
        size=12, bold=True, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE  5 — What we measured (image features)
# ═════════════════════════════════════════════════════════════════════════════

def slide_features(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Что именно мы измеряем на снимке",
           "20 признаков — числовые характеристики из пикселей МРТ")
    footer(s, 5)

    img(s, FIGURES / "image_analysis_panels.png",
        Inches(0.3), Inches(1.2), Inches(7.5), Inches(4.5))
    txt(s, "Один МРТ-срез превращается в таблицу из 20 чисел",
        Inches(0.3), Inches(5.75), Inches(7.5), Inches(0.4),
        size=13, color=C_DGRAY, italic=True)

    groups = [
        (C_BLUE,   "Яркость пикселей",  "среднее, разброс, форма распределения"),
        (C_GREEN,  "Резкость снимка",   "насколько чёткие границы тканей"),
        (C_ORANGE, "Текстура",          "однородность, контраст, энергия"),
        (C_BLUE,   "Геометрия",         "площадь кости, площадь хряща, щель"),
        (C_GREEN,  "Клинические",       "остеофиты, индекс склероза"),
    ]
    card(s, Inches(8.1), Inches(1.2), Inches(5.0), Inches(5.9))
    txt(s, "Группы признаков:",
        Inches(8.25), Inches(1.3), Inches(4.7), Inches(0.4),
        size=16, bold=True, color=C_LGRAY)
    y = Inches(1.85)
    for col, grp, desc in groups:
        rect(s, Inches(8.25), y, Inches(0.06), Inches(0.9), col)
        txt(s, grp, Inches(8.4), y + Inches(0.02), Inches(4.5), Inches(0.38),
            size=14, bold=True, color=col)
        txt(s, desc, Inches(8.4), y + Inches(0.42), Inches(4.5), Inches(0.42),
            size=12, color=C_LGRAY)
        y += Inches(1.02)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE  6 — Feature importance
# ═════════════════════════════════════════════════════════════════════════════

def slide_importance(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Какие признаки важнее всего?",
           "Random Forest показывает, что больше всего влияет на диагноз")
    footer(s, 6)

    img(s, FIGURES / "feature_importance.png",
        Inches(0.3), Inches(1.2), Inches(7.5), Inches(5.9))

    top5 = [
        (C_ORANGE, "#1", "Ширина суставной щели",   "Главный признак — щель сужается при артрите"),
        (C_BLUE,   "#2", "Остеофиты (наросты)",      "Костные шипы вокруг сустава"),
        (C_GREEN,  "#3", "Площадь хряща",            "Хрящ истончается по мере прогрессии"),
        (C_BLUE,   "#4", "Спектральная энергия",     "Изменение текстуры снимка"),
        (C_DGRAY,  "#5", "Резкость снимка",          "Размытие границ при артрите"),
    ]

    card(s, Inches(8.05), Inches(1.2), Inches(5.1), Inches(5.9))
    txt(s, "Топ-5 признаков",
        Inches(8.2), Inches(1.3), Inches(4.8), Inches(0.45),
        size=18, bold=True, color=C_LGRAY)
    y = Inches(1.9)
    for col, num, name, desc in top5:
        rect(s, Inches(8.2), y, Inches(0.55), Inches(0.9), col)
        txt(s, num, Inches(8.2), y + Inches(0.18), Inches(0.55), Inches(0.55),
            size=14, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
        txt(s, name, Inches(8.85), y + Inches(0.02), Inches(4.1), Inches(0.38),
            size=14, bold=True, color=col)
        txt(s, desc, Inches(8.85), y + Inches(0.44), Inches(4.1), Inches(0.4),
            size=12, color=C_LGRAY)
        y += Inches(1.0)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE  7 — EDA (3 different graphs)
# ═════════════════════════════════════════════════════════════════════════════

def slide_eda(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Что показал анализ данных")
    footer(s, 7)

    img(s, FIGURES / "class_distribution.png",
        Inches(0.3), Inches(1.15), Inches(4.3), Inches(3.0))
    img(s, FIGURES / "distribution_edge_density.png",
        Inches(4.75), Inches(1.15), Inches(4.3), Inches(3.0))
    img(s, FIGURES / "correlation_matrix.png",
        Inches(9.2), Inches(1.15), Inches(3.9), Inches(3.0))

    captions = [
        (Inches(0.3),  Inches(4.22), Inches(4.3),
         "Дисбаланс классов: снимков KL-2 и KL-3 больше, чем KL-0 и KL-4"),
        (Inches(4.75), Inches(4.22), Inches(4.3),
         "Резкость границ снижается при тяжёлых стадиях — ткани смешиваются"),
        (Inches(9.2),  Inches(4.22), Inches(3.9),
         "Некоторые признаки сильно коррелируют — нужна нелинейная модель"),
    ]
    for lx, ly, lw, cap in captions:
        txt(s, cap, lx, ly, lw, Inches(0.55),
            size=12, color=C_DGRAY, italic=True, align=PP_ALIGN.CENTER)

    findings = [
        (C_ORANGE, "График 1: дисбаланс классов — нужно учитывать при оценке точности"),
        (C_BLUE,   "График 2: резкость падает при тяжёлых стадиях — признак информативен"),
        (C_GREEN,  "График 3: сильные корреляции — линейная модель не оптимальна"),
    ]
    y = Inches(4.9)
    for col, line in findings:
        card(s, Inches(0.3), y, Inches(12.7), Inches(0.68), col)
        txt(s, line, Inches(0.55), y + Inches(0.1), Inches(12.2), Inches(0.5),
            size=15)
        y += Inches(0.76)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE  8 — PCA
# ═════════════════════════════════════════════════════════════════════════════

def slide_pca(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Визуализация данных (PCA)",
           "Можно ли разделить степени артрита одной линией?")
    footer(s, 8)

    img(s, FIGURES / "pca_projection.png",
        Inches(0.3), Inches(1.2), Inches(7.2), Inches(5.7))
    txt(s, "Каждая точка — один снимок. Цвет — степень артрита (KL 0–4)",
        Inches(0.3), Inches(6.98), Inches(7.2), Inches(0.4),
        size=11, color=C_DGRAY, italic=True)

    card(s, Inches(7.7), Inches(1.2), Inches(5.4), Inches(5.7))
    txt(s, "Как устроен этот график?",
        Inches(7.9), Inches(1.32), Inches(5.0), Inches(0.45),
        size=18, bold=True, color=C_BLUE)

    points = [
        (C_LGRAY,  "Таблица N×20 (N снимков, 20 параметров)"),
        (C_LGRAY,  "PCA строит 2 новые оси, сохраняя максимум информации"),
        (C_LGRAY,  "Каждый снимок становится точкой на 2D-графике"),
        (C_LGRAY,  "Это только для визуализации — сжимаем 20D → 2D"),
        (C_ORANGE, "Классы перемешаны — одну линию провести нельзя"),
        (C_GREEN,  "Random Forest строит нелинейные границы — справляется"),
    ]
    y = Inches(1.95)
    for col, line in points:
        if col == C_ORANGE:
            rect(s, Inches(7.9), y, Inches(5.0), Inches(1.0), C_CARD)
            rect(s, Inches(7.9), y, Inches(0.07), Inches(1.0), C_ORANGE)
        elif col == C_GREEN:
            rect(s, Inches(7.9), y, Inches(5.0), Inches(1.0), C_CARD)
            rect(s, Inches(7.9), y, Inches(0.07), Inches(1.0), C_GREEN)
        txt(s, line, Inches(8.07), y + Inches(0.12), Inches(4.8), Inches(0.78),
            size=14, color=col)
        y += Inches(1.12)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE  9 — Segmentation (Otsu vs U-Net vs nnU-Net)
# ═════════════════════════════════════════════════════════════════════════════

def slide_segmentation(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Сегментация тканей",
           "Это отдельная задача: выделение кости и хряща, НЕ определение стадии KL")
    footer(s, 9)

    # Pipeline steps
    steps = [
        (C_BLUE,   "Шаг 1\nВходной снимок",  "Серый МРТ-срез\nколенного сустава"),
        (C_BLUE,   "Шаг 2\nПорог Otsu",       "Автоматически\nотделяет ткань от фона"),
        (C_GREEN,  "Шаг 3\nМорфология",        "Заполняет\nпробелы и дыры"),
        (C_GREEN,  "Шаг 4\nКомпоненты",        "Находит отдельные\nобласти тканей"),
        (C_ORANGE, "Шаг 5\nКласс по размеру",  "Большая → кость\nМалая → хрящ"),
    ]
    bw = Inches(2.3)
    for i, (col, title, desc) in enumerate(steps):
        lx = Inches(0.3) + i * (bw + Inches(0.25))
        rect(s, lx, Inches(1.3), bw, Inches(3.2), C_CARD)
        rect(s, lx, Inches(1.3), bw, Inches(0.12), col)
        txt(s, title, lx, Inches(1.5), bw, Inches(0.9),
            size=15, bold=True, color=col, align=PP_ALIGN.CENTER)
        txt(s, desc, lx, Inches(2.5), bw, Inches(1.8),
            size=13, color=C_LGRAY, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            txt(s, "→", lx + bw + Inches(0.03), Inches(2.5),
                Inches(0.22), Inches(0.5),
                size=22, bold=True, color=C_DGRAY, align=PP_ALIGN.CENTER)

    # Methods comparison
    txt(s, "Сравнение методов сегментации (Dice = мера точности: 1.0 = идеально):",
        Inches(0.3), Inches(4.65), Inches(12.5), Inches(0.4),
        size=15, bold=True, color=C_LGRAY)

    methods = [
        (C_DGRAY, "Otsu\n(наш baseline)",  "Dice: 0.20–0.35\nПростой, быстрый"),
        (C_BLUE,  "U-Net\n(нейросеть)",    "Dice: 0.85–0.90\nТребует обучения"),
        (C_GREEN, "nnU-Net\n(лучший)",     "Dice: 0.90–0.95\nСостояние искусства"),
    ]
    for i, (col, name, metric) in enumerate(methods):
        lx = Inches(0.3) + i * Inches(4.25)
        rect(s, lx, Inches(5.15), Inches(4.0), Inches(2.1), C_CARD)
        rect(s, lx, Inches(5.15), Inches(4.0), Inches(0.12), col)
        txt(s, name, lx, Inches(5.32), Inches(4.0), Inches(0.85),
            size=18, bold=True, color=col, align=PP_ALIGN.CENTER)
        txt(s, metric, lx, Inches(6.2), Inches(4.0), Inches(0.95),
            size=14, color=C_LGRAY, align=PP_ALIGN.CENTER)

    txt(s, "Naш baseline (Otsu) — это точка отсчёта. U-Net и nnU-Net — план на развитие.",
        Inches(0.3), Inches(7.18), Inches(12.5), Inches(0.25),
        size=11, color=C_DGRAY, italic=True, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Classification: LR vs RF
# ═════════════════════════════════════════════════════════════════════════════

def slide_classification(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Классификация KL-стадии",
           "Logistic Regression = baseline (точка отсчёта)   |   Random Forest = наш финальный выбор")
    footer(s, 10)

    img(s, FIGURES / "confusion_matrices.png",
        Inches(0.3), Inches(1.2), Inches(6.8), Inches(4.3))
    txt(s, "Матрица ошибок: строки = реальный класс, столбцы = предсказание.\n"
           "Чем ярче диагональ — тем точнее модель.",
        Inches(0.3), Inches(5.55), Inches(6.8), Inches(0.65),
        size=12, color=C_DGRAY, italic=True)

    img(s, FIGURES / "roc_curves.png",
        Inches(7.2), Inches(1.2), Inches(5.9), Inches(4.3))
    txt(s, "ROC-кривая: чем выше линия — тем лучше модель различает классы",
        Inches(7.2), Inches(5.55), Inches(5.9), Inches(0.65),
        size=12, color=C_DGRAY, italic=True)

    models = [
        (C_DGRAY, "Logistic Regression",
         "Простая модель — используем как контрольную точку"),
        (C_GREEN, "Random Forest ✓",
         "Ансамбль деревьев — AUC и F1 выше по всем классам"),
    ]
    for i, (col, name, desc) in enumerate(models):
        lx = Inches(0.3) + i * Inches(6.4)
        card(s, lx, Inches(6.3), Inches(6.1), Inches(1.05), col)
        txt(s, name, lx + Inches(0.2), Inches(6.37), Inches(5.7), Inches(0.38),
            size=16, bold=True, color=col)
        txt(s, desc, lx + Inches(0.2), Inches(6.77), Inches(5.7), Inches(0.5),
            size=13, color=C_LGRAY)

    txt(s, "Сравнение с Logistic Regression нужно, чтобы показать: улучшение случайным образом не объяснить.",
        Inches(0.3), Inches(7.18), Inches(12.7), Inches(0.25),
        size=10, color=C_DGRAY, italic=True, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Metrics / Results
# ═════════════════════════════════════════════════════════════════════════════

def slide_results(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Итоговые метрики")
    footer(s, 11)

    img(s, FIGURES / "classification_metrics.png",
        Inches(0.3), Inches(1.2), Inches(8.5), Inches(5.9))

    card(s, Inches(9.0), Inches(1.2), Inches(4.1), Inches(5.9))
    txt(s, "Что означают метрики?",
        Inches(9.15), Inches(1.32), Inches(3.8), Inches(0.45),
        size=16, bold=True, color=C_LGRAY)

    defs = [
        (C_BLUE,   "Accuracy",  "Сколько % ответов верных"),
        (C_GREEN,  "Precision", "Из найденных — сколько правильных"),
        (C_ORANGE, "Recall",    "Из реальных — сколько нашли"),
        (C_BLUE,   "F1-score",  "Баланс между точностью и полнотой"),
        (C_GREEN,  "ROC-AUC",   "Площадь под кривой: 1.0 = идеал"),
    ]
    y = Inches(1.95)
    for col, name, desc in defs:
        rect(s, Inches(9.15), y, Inches(0.06), Inches(0.85), col)
        txt(s, name, Inches(9.3), y + Inches(0.02), Inches(3.6), Inches(0.35),
            size=13, bold=True, color=col)
        txt(s, desc, Inches(9.3), y + Inches(0.4), Inches(3.6), Inches(0.4),
            size=12, color=C_LGRAY)
        y += Inches(0.98)

    txt(s, "Если у простой модели Accuracy выглядит похожей — это из-за дисбаланса классов.\n"
           "Важнее F1 и AUC по всем классам — там Random Forest стабильнее.",
        Inches(0.3), Inches(7.08), Inches(12.7), Inches(0.35),
        size=10, color=C_DGRAY, italic=True, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Conclusions (exactly 3)
# ═════════════════════════════════════════════════════════════════════════════

def slide_conclusions(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Выводы")
    footer(s, 12)

    conclusions = [
        (C_BLUE,
         "1) KL-стадию можно предсказывать автоматически",
         "По МРТ-снимку модель определяет стадию артрита (KL-0 ... KL-4) без участия врача"),
        (C_GREEN,
         "2) 20 параметров вычисляются прямо из снимка",
         "Яркость, резкость, текстура, геометрия — никакого ручного ввода"),
        (C_ORANGE,
         "3) Random Forest — практичный и объяснимый выбор",
         "Нелинейные данные, высокий AUC, показывает важность признаков, устойчив к дисбалансу"),
    ]
    y = Inches(1.5)
    for col, title, desc in conclusions:
        card(s, Inches(0.3), y, Inches(12.7), Inches(1.3), col)
        txt(s, title, Inches(0.55), y + Inches(0.1),
            Inches(12.0), Inches(0.45), size=20, bold=True, color=col)
        txt(s, desc, Inches(0.55), y + Inches(0.63),
            Inches(12.0), Inches(0.55), size=16, color=C_WHITE)
        y += Inches(1.45)

    # "Next steps" mini-bar
    card(s, Inches(0.3), Inches(6.1), Inches(12.7), Inches(1.1))
    txt(s, "Следующие шаги:",
        Inches(0.5), Inches(6.18), Inches(2.5), Inches(0.38),
        size=15, bold=True, color=C_DGRAY)
    nexts = [
        (C_BLUE,   "Обучить U-Net\n(нейросеть)"),
        (C_GREEN,  "Проверить на\nреальных данных"),
        (C_ORANGE, "Интерфейс\nдля врача"),
    ]
    for i, (col, line) in enumerate(nexts):
        lx = Inches(2.8) + i * Inches(3.5)
        rect(s, lx, Inches(6.2), Inches(0.08), Inches(0.85), col)
        txt(s, line, lx + Inches(0.18), Inches(6.2),
            Inches(3.1), Inches(0.85), size=14, color=C_LGRAY)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Final / Thank you
# ═════════════════════════════════════════════════════════════════════════════

def slide_final(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    rect(s, 0, 0, SW, Inches(1.2), C_CARD)
    rect(s, 0, 0, Inches(0.12), SH, C_BLUE)
    rect(s, 0, SH - Inches(0.07), SW, Inches(0.07), C_BLUE)

    txt(s, "Спасибо за внимание!",
        Inches(1.0), Inches(1.7), Inches(11.0), Inches(1.1),
        size=44, bold=True, align=PP_ALIGN.CENTER)

    rect(s, Inches(3.5), Inches(3.0), Inches(6.2), Inches(0.06), C_GREEN)

    txt(s, "Маратов Ерасыл Балканович",
        Inches(1.0), Inches(3.2), Inches(11.0), Inches(0.5),
        size=20, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)

    txt(s, "7М06101  |  AIU  |  2025–2026",
        Inches(1.0), Inches(3.8), Inches(11.0), Inches(0.4),
        size=15, color=C_DGRAY, align=PP_ALIGN.CENTER)

    tags = ["МРТ + ML", "20 признаков", "Random Forest", "Dice / AUC", "U-Net — план"]
    tx = Inches(1.55)
    for tag in tags:
        w = Inches(2.1)
        rect(s, tx, Inches(4.9), w, Inches(0.55), C_BLUE)
        txt(s, tag, tx, Inches(4.95), w, Inches(0.45),
            size=14, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
        tx += w + Inches(0.15)

    txt(s, f"{TOTAL} / {TOTAL}",
        SW - Inches(1.0), SH - Inches(0.38), Inches(0.85), Inches(0.3),
        size=11, color=C_DGRAY, align=PP_ALIGN.RIGHT)


# ═════════════════════════════════════════════════════════════════════════════
# BUILD
# ═════════════════════════════════════════════════════════════════════════════

def build():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    slide_title(prs)
    slide_problem(prs)
    slide_data(prs)
    slide_params_and_rf(prs)
    slide_features(prs)
    slide_importance(prs)
    slide_eda(prs)
    slide_pca(prs)
    slide_segmentation(prs)
    slide_classification(prs)
    slide_results(prs)
    slide_conclusions(prs)
    slide_final(prs)

    out = Path(__file__).parent / "knee_mri_presentation_easy_explained.pptx"
    prs.save(str(out))
    print(f"Готово: {out}  ({len(prs.slides)} слайдов)")


if __name__ == "__main__":
    build()
