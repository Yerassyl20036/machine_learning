#!/usr/bin/env python3
"""
Knee MRI Segmentation — PowerPoint (in Russian), SEGMENTATION topic.
Topic: Методы автоматической сегментации костных и хрящевых тканей
       коленного сустава по данным МРТ

Usage : python generate_slides.py
Output: knee_mri_presentation_easy_explained.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── Paths ────────────────────────────────────────────────────────────────────
SEG_FIGS = Path(__file__).parent / "results" / "seg_figures"

# ─── Colors ───────────────────────────────────────────────────────────────────
C_BG    = RGBColor(0x0E, 0x1A, 0x2B)
C_CARD  = RGBColor(0x12, 0x22, 0x38)
C_DARK  = RGBColor(0x1C, 0x33, 0x4D)
C_BLUE  = RGBColor(0x00, 0x99, 0xFF)
C_GREEN = RGBColor(0x00, 0xCC, 0x99)
C_ORANGE= RGBColor(0xFF, 0x99, 0x00)
C_RED   = RGBColor(0xFF, 0x44, 0x55)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LGRAY = RGBColor(0xCC, 0xDD, 0xEE)
C_DGRAY = RGBColor(0x77, 0x99, 0xBB)

SW = Inches(13.33)
SH = Inches(7.5)
TOTAL = 15


# ═══════════════════════════════════════════════════════════════════════════════
# Primitive helpers
# ═══════════════════════════════════════════════════════════════════════════════

def bg(slide, color=C_BG):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def rect(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def txt(slide, content, l, t, w, h, size=18, bold=False,
        color=C_WHITE, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = content
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def bullets(slide, lines, l, t, w, h, size=16, color=C_LGRAY):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(7)
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = "●  " + line
        r.font.size = Pt(size)
        r.font.color.rgb = color


def img(slide, path, l, t, w, h):
    p = Path(path)
    if p.exists():
        slide.shapes.add_picture(str(p), l, t, w, h)
    else:
        rect(slide, l, t, w, h, C_DARK)
        txt(slide, f"[ {p.name} ]", l, t + h // 3, w, Inches(0.4),
            size=11, color=C_DGRAY, align=PP_ALIGN.CENTER)


def header(slide, title, sub=None):
    rect(slide, 0, 0, SW, Inches(1.05), C_DARK)
    rect(slide, 0, 0, Inches(0.1), Inches(1.05), C_BLUE)
    txt(slide, title, Inches(0.25), Inches(0.1), Inches(12.5), Inches(0.6),
        size=26, bold=True)
    if sub:
        txt(slide, sub, Inches(0.25), Inches(0.68), Inches(12.5), Inches(0.35),
            size=13, color=C_BLUE, italic=True)


def footer(slide, num):
    rect(slide, 0, SH - Inches(0.07), SW, Inches(0.07), C_BLUE)
    txt(slide, f"{num} / {TOTAL}",
        SW - Inches(1.0), SH - Inches(0.38), Inches(0.85), Inches(0.3),
        size=11, color=C_DGRAY, align=PP_ALIGN.RIGHT)


def card(slide, l, t, w, h, accent=None):
    rect(slide, l, t, w, h, C_DARK)
    if accent:
        rect(slide, l, t, Inches(0.07), h, accent)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    rect(s, 0, 0, Inches(7.2), SH, C_DARK)
    rect(s, 0, 0, Inches(0.12), SH, C_BLUE)

    txt(s, "Методы автоматической\nсегментации костных\nи хрящевых тканей\nколенного сустава",
        Inches(0.3), Inches(0.6), Inches(6.6), Inches(3.0), size=28, bold=True)
    txt(s, "по данным МРТ",
        Inches(0.3), Inches(3.7), Inches(6.6), Inches(0.55),
        size=20, color=C_BLUE, bold=True)
    rect(s, Inches(0.3), Inches(4.5), Inches(5.8), Inches(0.05), C_GREEN)
    txt(s, "Маратов Ерасыл Балканович",
        Inches(0.3), Inches(4.65), Inches(6.3), Inches(0.45),
        size=18, bold=True, color=C_GREEN)
    txt(s, "7М06101  |  AIU  |  2025–2026  |  Алматы",
        Inches(0.3), Inches(5.2), Inches(6.3), Inches(0.4),
        size=13, color=C_DGRAY)
    img(s, SEG_FIGS / "kl_seg_samples.png",
        Inches(7.3), Inches(0.1), Inches(5.9), Inches(7.3))
    footer(s, 1)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem / Motivation
# ═══════════════════════════════════════════════════════════════════════════════

def slide_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Зачем нужна автоматическая сегментация?")
    footer(s, 2)

    card(s, Inches(0.3), Inches(1.2), Inches(12.7), Inches(1.3), C_ORANGE)
    txt(s, "Остеоартрит коленного сустава — 520 млн пациентов в мире. Ранняя диагностика критична.",
        Inches(0.55), Inches(1.35), Inches(12.2), Inches(0.85), size=19, bold=True)

    card(s, Inches(0.3), Inches(2.65), Inches(6.1), Inches(4.4), C_ORANGE)
    txt(s, "Сегодня — вручную",
        Inches(0.55), Inches(2.75), Inches(5.7), Inches(0.45),
        size=18, bold=True, color=C_ORANGE)
    bullets(s, [
        "Врач обводит кость и хрящ на каждом МРТ-срезе",
        "30–90 минут на одного пациента",
        "Разные врачи — разные оценки одного снимка",
        "МРТ даёт 20–40 срезов → сотни часов работы",
        "Поздняя диагностика ухудшает прогноз",
    ], Inches(0.55), Inches(3.3), Inches(5.7), Inches(3.5), size=15)

    card(s, Inches(6.6), Inches(2.65), Inches(6.4), Inches(4.4), C_GREEN)
    txt(s, "Цель этой работы — автоматически",
        Inches(6.85), Inches(2.75), Inches(6.0), Inches(0.45),
        size=18, bold=True, color=C_GREEN)
    bullets(s, [
        "Алгоритм находит и выделяет кость и хрящ без врача",
        "Секунды вместо часа на любом снимке",
        "Воспроизводимый результат — без субъективности",
        "Основа для диагностики стадии артрита (шкала KL)",
        "Сравниваем методы: Otsu → Double-Otsu → U-Net",
    ], Inches(6.85), Inches(3.3), Inches(6.0), Inches(3.5), size=15)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Dataset
# ═══════════════════════════════════════════════════════════════════════════════

def slide_dataset(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Данные", "Датасет: 3D Knee MRI Cartilage Segmentation (OAI / Kaggle)")
    footer(s, 3)

    stats = [
        (C_BLUE,   "Источник",   "Osteoarthritis Initiative\n(OAI, Kaggle)"),
        (C_GREEN,  "Объём",      "~1733 2D-срезов МРТ\nс экспертными масками"),
        (C_ORANGE, "Стадии KL",  "5 степеней артрита\nKL-0 … KL-4"),
        (C_RED,    "Оценка",     "Синтетические данные\n(Kaggle: ручная загрузка)"),
    ]
    for i, (col, label, val) in enumerate(stats):
        lx = Inches(0.3) + i * Inches(3.25)
        rect(s, lx, Inches(1.2), Inches(3.0), Inches(1.7), C_DARK)
        rect(s, lx, Inches(1.2), Inches(3.0), Inches(0.1), col)
        txt(s, label, lx, Inches(1.35), Inches(3.0), Inches(0.4),
            size=13, color=col, align=PP_ALIGN.CENTER)
        txt(s, val, lx, Inches(1.72), Inches(3.0), Inches(1.1),
            size=14, bold=True, align=PP_ALIGN.CENTER)

    img(s, SEG_FIGS / "dataset_split.png",
        Inches(0.3), Inches(3.05), Inches(5.5), Inches(4.1))

    kl_rows = [
        (C_GREEN,                       "KL-0", "Норма — сустав здоровый"),
        (C_BLUE,                        "KL-1", "Ранние изменения"),
        (C_ORANGE,                      "KL-2", "Умеренная стадия"),
        (RGBColor(0xFF, 0x66, 0x33),    "KL-3", "Выраженная стадия"),
        (C_RED,                         "KL-4", "Тяжёлая стадия"),
    ]
    card(s, Inches(6.1), Inches(3.05), Inches(7.0), Inches(3.5))
    txt(s, "Шкала KL (Kellgren–Lawrence):",
        Inches(6.35), Inches(3.15), Inches(6.5), Inches(0.4),
        size=15, bold=True, color=C_LGRAY)
    y = Inches(3.65)
    for col, stage, desc in kl_rows:
        rect(s, Inches(6.35), y, Inches(0.65), Inches(0.44), col)
        txt(s, stage, Inches(6.35), y + Inches(0.04), Inches(0.65), Inches(0.36),
            size=12, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
        txt(s, desc, Inches(7.1), y + Inches(0.04), Inches(5.8), Inches(0.38),
            size=14, color=C_LGRAY)
        y += Inches(0.58)

    card(s, Inches(6.1), Inches(6.65), Inches(7.0), Inches(0.55), C_ORANGE)
    txt(s, "Датасет содержит экспертные маски (ground truth): каждый пиксель помечен "
           "радиологом как кость, хрящ или фон. Именно с ними мы сравниваем наш алгоритм (Dice).",
        Inches(6.25), Inches(6.7), Inches(6.7), Inches(0.45),
        size=10, color=C_WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — What is segmentation + Dice
# ═══════════════════════════════════════════════════════════════════════════════

def slide_what_is_seg(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Что такое сегментация и как измерить точность?",
           "Задача: каждому пикселю присвоить метку — фон / кость / хрящ")
    footer(s, 4)

    img(s, SEG_FIGS / "dice_explained.png",
        Inches(0.3), Inches(1.2), Inches(8.0), Inches(5.9))

    card(s, Inches(8.55), Inches(1.2), Inches(4.5), Inches(5.9))
    txt(s, "Метрика Dice Score",
        Inches(8.75), Inches(1.32), Inches(4.1), Inches(0.45),
        size=17, bold=True, color=C_ORANGE)
    txt(s, "Dice = 2×TP / (2×TP + FP + FN)",
        Inches(8.75), Inches(1.85), Inches(4.1), Inches(0.5),
        size=14, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

    scale = [
        (C_RED,    "0.0–0.3", "Плохо\n(Otsu v1)"),
        (C_ORANGE, "0.3–0.6", "Удовл."),
        (C_BLUE,   "0.6–0.8", "Хорошо\n(наш улучш.)"),
        (C_GREEN,  "0.8–1.0", "Отлично\n(U-Net)"),
    ]
    y = Inches(2.55)
    for col, rng, lbl in scale:
        rect(s, Inches(8.75), y, Inches(0.9), Inches(0.9), col)
        txt(s, rng, Inches(8.75), y + Inches(0.2), Inches(0.9), Inches(0.5),
            size=11, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
        txt(s, lbl, Inches(9.75), y + Inches(0.15), Inches(3.1), Inches(0.6),
            size=12, color=C_LGRAY)
        y += Inches(1.0)

    txt(s, "TP = пиксели, верно найденные алгоритмом\n"
           "FP = лишние (алгоритм нашёл, но их нет)\n"
           "FN = пропущенные (алгоритм не нашёл)",
        Inches(8.75), Inches(6.65), Inches(4.1), Inches(0.55),
        size=10, color=C_DGRAY, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — First approach: simple Otsu (brightness only)
# ═══════════════════════════════════════════════════════════════════════════════

def slide_first_approach(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Первый подход: простой порог Otsu",
           "Классификация только по яркости пикселя — без учёта соседей и положения")
    footer(s, 5)

    img(s, SEG_FIGS / "seg_pipeline.png",
        Inches(0.3), Inches(1.15), Inches(12.7), Inches(4.3))

    steps = [
        (C_DGRAY, "Шаг 1: Нормализация",
                  "Привести все пиксели к диапазону 0–255"),
        (C_ORANGE,"Шаг 2: Порог Otsu (один)",
                  "Автоматически найти t₁ по гистограмме.\nВыше t₁ = ткань. Ниже = фон."),
        (C_RED,   "Шаг 3: Самый большой блоб = кость",
                  "Связные компоненты: наибольшая область → кость,\nостальные → хрящ."),
    ]
    y = Inches(5.6)
    for col, title, desc in steps:
        card(s, Inches(0.3), y, Inches(12.7), Inches(0.62), col)
        txt(s, title, Inches(0.55), y + Inches(0.05), Inches(3.8), Inches(0.52),
            size=13, bold=True, color=col)
        txt(s, desc, Inches(4.4), y + Inches(0.07), Inches(9.0), Inches(0.48),
            size=12, color=C_WHITE)
        y += Inches(0.69)

    card(s, Inches(0.3), Inches(7.12), Inches(12.7), Inches(0.3), C_RED)
    txt(s, "Ключевое ограничение: один порог делит только ТКАНЬ / ФОН. Кость и хрящ попадают в один класс.",
        Inches(0.5), Inches(7.16), Inches(12.3), Inches(0.24),
        size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — First results: Dice 0.20–0.35, cartilage = 0
# ═══════════════════════════════════════════════════════════════════════════════

def slide_first_results(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Результаты первого подхода",
           "Otsu с одним порогом — тест на синтетических данных (KL-0 … KL-4)")
    footer(s, 6)

    img(s, SEG_FIGS / "seg_samples.png",
        Inches(0.3), Inches(1.2), Inches(7.8), Inches(4.0))

    results_v1 = [
        (C_DGRAY,  "Dice — кость",  "0.20–0.35", "Низкий, но кость хотя бы находится"),
        (C_RED,    "Dice — хрящ",   "≈ 0.000",   "Хрящ полностью не выделяется!"),
        (C_RED,    "Хрящевых px",   "~0–180",    "Алгоритм почти не находит хрящ"),
        (C_ORANGE, "Скорость",      "< 0.1 сек", "Единственный плюс — очень быстро"),
    ]
    y = Inches(1.2)
    for col, label, val, note in results_v1:
        card(s, Inches(8.1), y, Inches(5.0), Inches(0.82), col)
        txt(s, label, Inches(8.3), y + Inches(0.05), Inches(2.0), Inches(0.36),
            size=13, color=col, bold=True)
        txt(s, val, Inches(10.4), y + Inches(0.1), Inches(1.5), Inches(0.55),
            size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        txt(s, note, Inches(8.3), y + Inches(0.46), Inches(4.6), Inches(0.3),
            size=11, color=C_LGRAY)
        y += Inches(0.9)

    card(s, Inches(8.1), Inches(5.2), Inches(5.0), Inches(2.0), C_RED)
    txt(s, "Почему хрящ не находится?",
        Inches(8.3), Inches(5.3), Inches(4.6), Inches(0.42),
        size=15, bold=True, color=C_RED)
    bullets(s, [
        "Один порог делит только ткань / фон",
        "Кость и хрящ оба выше порога → один блоб",
        "Самый большой блоб → всё становится «костью»",
        "Хрящ получает ноль пикселей → Dice = 0",
    ], Inches(8.3), Inches(5.78), Inches(4.6), Inches(1.3), size=12)

    card(s, Inches(0.3), Inches(5.35), Inches(7.8), Inches(1.85))
    txt(s, "Вывод: первый метод провалился для хряща.\nНужно найти другой способ разделить кость и хрящ.",
        Inches(0.5), Inches(5.45), Inches(7.4), Inches(1.6),
        size=16, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Why Otsu fails — analysis
# ═══════════════════════════════════════════════════════════════════════════════

def slide_why_fails(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Анализ провала: почему один порог не работает?",
           "Физика МРТ-сигнала + ограничения порогового метода")
    footer(s, 7)

    img(s, SEG_FIGS / "dice_explained.png",
        Inches(0.3), Inches(1.2), Inches(5.8), Inches(5.9))

    reasons = [
        (C_RED,    "Одинаковая яркость",
                   "В МРТ кость и хрящ дают схожий сигнал. "
                   "Один порог t₁ помещает оба в один класс «ткань»."),
        (C_ORANGE, "Хрящ — тонкий слой",
                   "Толщина хряща 2–5 мм. Он расположен вплотную к кости. "
                   "Алгоритм не видит, где заканчивается кость и начинается хрящ."),
        (C_BLUE,   "Нет знания о форме",
                   "Otsu не знает: «кость — большая эллипсовидная масса, "
                   "хрящ — тонкая дуга рядом с ней». Решение только по яркости."),
        (C_GREEN,  "Что нужно сделать?",
                   "Найти ВТОРОЙ порог внутри ткани: разделить кость (средняя яркость) "
                   "от хряща (более высокая яркость в МРТ Т2)."),
    ]
    y = Inches(1.2)
    for col, title, desc in reasons:
        card(s, Inches(6.35), y, Inches(6.7), Inches(1.35), col)
        txt(s, title, Inches(6.55), y + Inches(0.08), Inches(6.3), Inches(0.42),
            size=15, bold=True, color=col)
        txt(s, desc, Inches(6.55), y + Inches(0.56), Inches(6.3), Inches(0.68),
            size=12, color=C_LGRAY)
        y += Inches(1.47)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Literature review: what techniques improve results
# ═══════════════════════════════════════════════════════════════════════════════

def slide_literature(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Обзор литературы: что улучшает сегментацию МРТ?",
           "Изучили статьи по сегментации коленного сустава → нашли ключевые техники")
    footer(s, 8)

    techniques = [
        (C_ORANGE, "Multi-level Otsu\n(Otsu, 1979; Liao, 2001)",
                   "Вместо одного порога — несколько.\n"
                   "Разделить пиксели на 3+ кластера по яркости.\n"
                   "→ Отдельный порог для кости и хряща.",
                   "Otsu IEEE Trans. 1979\nLiao et al. 2001"),
        (C_BLUE,   "CLAHE — контрастирование\n(Zuiderveld, 1994)",
                   "Адаптивное усиление локального контраста.\n"
                   "МРТ-снимки имеют неравномерную яркость.\n"
                   "→ Хрящ становится заметнее перед порогом.",
                   "Zuiderveld, Graphics Gems 1994"),
        (C_GREEN,  "Морфологические операции\n(Serra, 1982)",
                   "Closing (закрытие) заполняет дыры в маске.\n"
                   "Opening (открытие) удаляет шум.\n"
                   "→ Более чистые контуры кости и хряща.",
                   "Serra, Image Analysis\nand Math. Morph. 1982"),
        (C_DGRAY,  "U-Net / nnU-Net\n(Ronneberger 2015; Isensee 2021)",
                   "Свёрточная нейросеть. Учится на парах\n"
                   "(МРТ, маска). Dice 0.85–0.95.\n"
                   "→ Требует GPU и 1000+ примеров.",
                   "Ronneberger et al. MICCAI 2015\nIsensee et al. Nature Methods 2021"),
    ]
    for i, (col, title, desc, ref) in enumerate(techniques):
        lx = Inches(0.3) + (i % 2) * Inches(6.5)
        ty = Inches(1.2) + (i // 2) * Inches(3.0)
        card(s, lx, ty, Inches(6.2), Inches(2.75), col)
        txt(s, title, lx + Inches(0.2), ty + Inches(0.1), Inches(5.8), Inches(0.7),
            size=15, bold=True, color=col)
        txt(s, desc, lx + Inches(0.2), ty + Inches(0.85), Inches(5.8), Inches(1.3),
            size=12, color=C_WHITE)
        txt(s, ref, lx + Inches(0.2), ty + Inches(2.38), Inches(5.8), Inches(0.32),
            size=9, color=C_DGRAY, italic=True)

    card(s, Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.3), C_BLUE)
    txt(s, "Применили: Multi-level Otsu + морфология → Double-Otsu метод",
        Inches(0.5), Inches(7.14), Inches(12.3), Inches(0.24),
        size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Improved method: Double-Otsu
# ═══════════════════════════════════════════════════════════════════════════════

def slide_improved_method(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Улучшенный метод: Double-Otsu сегментация",
           "Два последовательных порога + морфологическое уточнение")
    footer(s, 9)

    img(s, SEG_FIGS / "seg_pipeline.png",
        Inches(0.3), Inches(1.15), Inches(12.7), Inches(3.8))

    steps = [
        (C_DGRAY,  "Шаг 1",  "Нормализация",
                   "Привести пиксели 0–255"),
        (C_ORANGE, "Шаг 2",  "Otsu t₁",
                   "Ткань vs фон. t₁ авто по всей гистограмме"),
        (C_BLUE,   "Шаг 3",  "Otsu t₂ внутри ткани",
                   "Только пиксели выше t₁. Находим t₂:\n< t₂ = кость, ≥ t₂ = хрящ"),
        (C_GREEN,  "Шаг 4",  "Морфология",
                   "Closing 5×5 → заполнить дыры\nOpening 3×3 → убрать шум"),
    ]
    y = Inches(5.2)
    w = Inches(3.1)
    for i, (col, num, title, desc) in enumerate(steps):
        lx = Inches(0.3) + i * Inches(3.25)
        card(s, lx, y, w, Inches(2.1), col)
        txt(s, num, lx + Inches(0.12), y + Inches(0.08), w, Inches(0.38),
            size=11, bold=True, color=col)
        txt(s, title, lx + Inches(0.12), y + Inches(0.45), w - Inches(0.2), Inches(0.45),
            size=14, bold=True, color=C_WHITE)
        txt(s, desc, lx + Inches(0.12), y + Inches(0.98), w - Inches(0.2), Inches(1.0),
            size=11, color=C_LGRAY)

    txt(s, "Ключевое отличие от первого метода: t₂ делает второй разрез ВНУТРИ ткани — отделяет кость (средняя яркость) от хряща (более высокая яркость в МРТ)",
        Inches(0.3), Inches(7.12), Inches(12.7), Inches(0.3),
        size=11, color=C_ORANGE, bold=True, italic=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — New results
# ═══════════════════════════════════════════════════════════════════════════════

def slide_new_results(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Результаты после улучшения",
           "Double-Otsu vs Simple Otsu — синтетические данные KL-0 … KL-4")
    footer(s, 10)

    img(s, SEG_FIGS / "dice_comparison.png",
        Inches(0.3), Inches(1.2), Inches(7.5), Inches(5.9))

    card(s, Inches(8.0), Inches(1.2), Inches(5.1), Inches(0.55))
    txt(s, "KL-стадия    Кость (v2)    Хрящ (v2)",
        Inches(8.1), Inches(1.28), Inches(4.9), Inches(0.38),
        size=12, bold=True, color=C_DGRAY)

    rows = [
        ("KL-0 (Норма)",     "0.93", "0.38", C_GREEN),
        ("KL-1 (Ранняя)",    "0.96", "0.27", C_BLUE),
        ("KL-2 (Умеренная)", "0.80", "0.29", C_ORANGE),
        ("KL-3 (Выражен.)",  "0.70", "0.10", RGBColor(0xFF, 0x66, 0x33)),
        ("KL-4 (Тяжёлая)",   "0.78", "0.05", C_RED),
        ("Среднее",          "0.83", "0.22", C_DGRAY),
    ]
    y = Inches(1.85)
    for label, bone, cart, col in rows:
        rect(s, Inches(8.0), y, Inches(5.1), Inches(0.5), C_DARK)
        rect(s, Inches(8.0), y, Inches(0.07), Inches(0.5), col)
        txt(s, label, Inches(8.15), y + Inches(0.07), Inches(2.3), Inches(0.36),
            size=12, color=C_LGRAY)
        txt(s, bone, Inches(10.5), y + Inches(0.05), Inches(1.0), Inches(0.4),
            size=14, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
        txt(s, cart, Inches(11.6), y + Inches(0.05), Inches(1.2), Inches(0.4),
            size=14, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)
        y += Inches(0.54)

    before_after = [
        (C_RED,   "Simple Otsu (v1)", "Кость: 0.20–0.35\nХрящ:  ≈ 0.000"),
        (C_GREEN, "Double-Otsu (v2)", "Кость: 0.68–0.96\nХрящ:  0.05–0.38"),
    ]
    y = Inches(5.25)
    for col, label, scores in before_after:
        card(s, Inches(8.0), y, Inches(5.1), Inches(0.85), col)
        txt(s, label, Inches(8.2), y + Inches(0.06), Inches(2.5), Inches(0.4),
            size=14, bold=True, color=col)
        txt(s, scores, Inches(10.8), y + Inches(0.1), Inches(2.1), Inches(0.65),
            size=13, bold=True, color=C_WHITE)
        y += Inches(0.92)

    txt(s, "Хрящ теперь находится! Для KL-4 Dice хряща низкий (0.05) — хрящ почти исчез.",
        Inches(8.0), Inches(7.12), Inches(5.1), Inches(0.3),
        size=10, color=C_DGRAY, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Why CV not DL — method choice
# ═══════════════════════════════════════════════════════════════════════════════

def slide_method_choice(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Почему классический CV, а не нейросеть?",
           "Обоснование выбора метода — Double-Otsu vs U-Net")
    footer(s, 11)

    card(s, Inches(0.3), Inches(1.2), Inches(6.1), Inches(5.9), C_DGRAY)
    txt(s, "Double-Otsu (наш выбор)",
        Inches(0.55), Inches(1.32), Inches(5.7), Inches(0.45),
        size=18, bold=True, color=C_DGRAY)
    bullets(s, [
        "НЕ требует обучения — нет нужды в 1000+ масках",
        "НЕ требует GPU — работает на любом компьютере",
        "Полностью интерпретируем: 2 числа (t₁, t₂) объясняют всё",
        "Скорость: < 0.1 секунды на снимок",
        "Базируется на физике МРТ (яркость = водный сигнал)",
        "Подходит для демо-системы и baseline оценки",
    ], Inches(0.55), Inches(1.88), Inches(5.7), Inches(4.5), size=13)

    card(s, Inches(6.6), Inches(1.2), Inches(6.4), Inches(5.9), C_BLUE)
    txt(s, "U-Net (следующий шаг)",
        Inches(6.85), Inches(1.32), Inches(6.0), Inches(0.45),
        size=18, bold=True, color=C_BLUE)
    bullets(s, [
        "Dice 0.85–0.95 — в 3-4 раза лучше нашего",
        "Учитывает форму, текстуру, контекст всего снимка",
        "Encoder-Decoder архитектура со skip-соединениями",
        "Требует GPU (NVIDIA ≥ 8GB) для обучения",
        "Требует 1000+ аннотированных МРТ-срезов",
        "2–7 дней обучения на реальных данных",
        "Augmentation: flip, rotation, Gaussian noise, elastic deform",
    ], Inches(6.85), Inches(1.88), Inches(6.0), Inches(4.5), size=13)

    card(s, Inches(0.3), Inches(7.12), Inches(12.7), Inches(0.3), C_ORANGE)
    txt(s, "Вывод: Double-Otsu — честный, прозрачный baseline. U-Net — цель для следующего этапа.",
        Inches(0.5), Inches(7.16), Inches(12.3), Inches(0.24),
        size=12, bold=True, color=C_BG, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Full comparison table Otsu v1 → v2 → U-Net
# ═══════════════════════════════════════════════════════════════════════════════

def slide_full_comparison(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Итоговое сравнение трёх методов")
    footer(s, 12)

    img(s, SEG_FIGS / "method_comparison.png",
        Inches(0.3), Inches(1.2), Inches(7.5), Inches(4.5))

    cols_h = ["Метод", "Dice кость", "Dice хрящ", "GPU?", "Обучение?", "Скорость"]
    rows_d = [
        ("Simple Otsu (v1)",  "0.20–0.35", "≈ 0.00",   "Нет", "Нет",    "< 0.1 с",  C_RED),
        ("Double-Otsu (v2)",  "0.68–0.96", "0.05–0.38","Нет", "Нет",    "< 0.1 с",  C_ORANGE),
        ("U-Net (лит-ра)",    "0.85–0.90", "0.80–0.87","Да",  "~7 дней","0.3 с",    C_BLUE),
        ("nnU-Net (лит-ра)",  "0.90–0.95", "0.88–0.93","Да",  "~7 дней","0.5 с",    C_GREEN),
    ]

    hx = Inches(8.0)
    hy = Inches(1.2)
    col_ws = [Inches(2.3), Inches(1.15), Inches(1.15), Inches(0.65), Inches(1.1), Inches(0.95)]
    for j, hdr in enumerate(cols_h):
        rect(s, hx + sum(col_ws[:j]), hy, col_ws[j], Inches(0.38), C_DARK)
        txt(s, hdr, hx + sum(col_ws[:j]) + Inches(0.05), hy + Inches(0.04),
            col_ws[j], Inches(0.3), size=10, bold=True, color=C_DGRAY)

    y = hy + Inches(0.4)
    for method, bone, cart, gpu, train, speed, col in rows_d:
        vals = [method, bone, cart, gpu, train, speed]
        for j, val in enumerate(vals):
            rect(s, hx + sum(col_ws[:j]), y, col_ws[j], Inches(0.55), C_CARD)
            rect(s, hx + sum(col_ws[:j]), y, Inches(0.06), Inches(0.55), col)
            c = col if j == 0 else C_WHITE
            txt(s, val, hx + sum(col_ws[:j]) + Inches(0.1), y + Inches(0.1),
                col_ws[j] - Inches(0.1), Inches(0.35), size=11, color=c, bold=(j==0))
        y += Inches(0.58)

    card(s, Inches(0.3), Inches(5.85), Inches(7.5), Inches(1.28))
    txt(s, "Что видно на изображении:",
        Inches(0.5), Inches(5.95), Inches(7.0), Inches(0.38),
        size=14, bold=True, color=C_LGRAY)
    bullets(s, [
        "Слева: исходный МРТ-снимок",
        "Центр: наш Double-Otsu (синий=кость, зелёный=хрящ)",
        "Справа: симулированный U-Net (для сравнения)",
    ], Inches(0.5), Inches(6.4), Inches(7.0), Inches(0.68), size=12)

    txt(s, "Источники: Ronneberger et al. 2015 (U-Net); Isensee et al. 2021 (nnU-Net)",
        Inches(8.0), Inches(7.12), Inches(5.1), Inches(0.3),
        size=9, color=C_DGRAY, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Web demo
# ═══════════════════════════════════════════════════════════════════════════════

def slide_demo(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Демонстрация: веб-интерфейс сегментации",
           "Загрузи МРТ-снимок → получи выделенную кость и хрящ + пошаговая визуализация")
    footer(s, 13)

    txt(s, "http://127.0.0.1:8000",
        Inches(0.3), Inches(1.35), Inches(12.7), Inches(0.7),
        size=26, bold=True, color=C_BLUE, align=PP_ALIGN.CENTER)

    features = [
        (C_DGRAY,  "5 примеров",    "Готовые снимки KL-0…KL-4\nнажми → сегментация сразу"),
        (C_ORANGE, "Drag & Drop",   "Загрузи свой PNG/JPG\nлюбой размер"),
        (C_BLUE,   "4 шага",        "Нормализация → Otsu t₁\n→ Otsu t₂ → морфология"),
        (C_GREEN,  "Результат",     "Синий=кость, зелёный=хрящ\n% площади + Dice"),
    ]
    y = Inches(2.25)
    for i, (col, title, desc) in enumerate(features):
        lx = Inches(0.3) + i * Inches(3.25)
        rect(s, lx, y, Inches(3.0), Inches(2.2), C_DARK)
        rect(s, lx, y, Inches(3.0), Inches(0.1), col)
        txt(s, title, lx, y + Inches(0.18), Inches(3.0), Inches(0.42),
            size=16, bold=True, color=col, align=PP_ALIGN.CENTER)
        txt(s, desc, lx + Inches(0.15), y + Inches(0.72), Inches(2.7), Inches(1.3),
            size=13, color=C_LGRAY, align=PP_ALIGN.CENTER)

    card(s, Inches(0.3), Inches(4.6), Inches(12.7), Inches(2.65))
    txt(s, "Пошаговая визуализация алгоритма в веб-интерфейсе:",
        Inches(0.55), Inches(4.7), Inches(12.0), Inches(0.38),
        size=14, bold=True, color=C_LGRAY)
    step_info = [
        (C_DGRAY,  "Шаг 1\nНорм.", "0–255"),
        (C_ORANGE, "Шаг 2\nOtsu t₁", "ткань/фон"),
        (C_BLUE,   "Шаг 3\nOtsu t₂", "кость/хрящ"),
        (C_GREEN,  "Шаг 4\nИтог", "overlay"),
    ]
    sx = Inches(0.55)
    for col, step, sub in step_info:
        rect(s, sx, Inches(5.15), Inches(2.9), Inches(1.95), C_CARD)
        rect(s, sx, Inches(5.15), Inches(2.9), Inches(0.08), col)
        txt(s, step, sx + Inches(0.1), Inches(5.28), Inches(2.7), Inches(0.65),
            size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
        txt(s, sub, sx + Inches(0.1), Inches(5.98), Inches(2.7), Inches(0.4),
            size=11, color=C_DGRAY, align=PP_ALIGN.CENTER)
        sx += Inches(3.15)

    txt(s, "Стек: FastAPI + OpenCV + Vanilla JS  |  Без GPU  |  < 0.1 сек",
        Inches(0.3), Inches(7.13), Inches(12.7), Inches(0.28),
        size=11, color=C_DGRAY, italic=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Conclusions
# ═══════════════════════════════════════════════════════════════════════════════

def slide_conclusions(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    header(s, "Выводы")
    footer(s, 14)

    conclusions = [
        (C_RED,
         "1) Первый метод (Simple Otsu) — хрящ не выделяется",
         "Один порог делит только ткань/фон. "
         "Кость и хрящ сливаются в один блоб → Dice хряща ≈ 0.000."),
        (C_ORANGE,
         "2) После изучения литературы — Double-Otsu с морфологией",
         "Два последовательных порога: t₁ (ткань/фон), t₂ (кость/хрящ внутри ткани). "
         "Dice кость: 0.68–0.96, Dice хрящ: 0.05–0.38."),
        (C_BLUE,
         "3) U-Net (следующий шаг) даёт Dice 0.85–0.95",
         "Нейросеть учится на форме и контексте, а не только на яркости. "
         "Требует GPU и аннотированных данных. Цель для развития работы."),
        (C_GREEN,
         "4) Точная сегментация → автоматический диагноз стадии артрита",
         "По площади хряща и ширине суставной щели можно автоматически "
         "определить KL-стадию остеоартрита без участия врача."),
    ]
    y = Inches(1.3)
    for col, title, desc in conclusions:
        card(s, Inches(0.3), y, Inches(12.7), Inches(1.2), col)
        txt(s, title, Inches(0.55), y + Inches(0.08),
            Inches(12.0), Inches(0.42), size=17, bold=True, color=col)
        txt(s, desc, Inches(0.55), y + Inches(0.58),
            Inches(12.0), Inches(0.52), size=13, color=C_WHITE)
        y += Inches(1.3)

    card(s, Inches(0.3), Inches(6.5), Inches(12.7), Inches(0.88))
    txt(s, "Следующие шаги:",
        Inches(0.5), Inches(6.56), Inches(2.0), Inches(0.35),
        size=13, bold=True, color=C_DGRAY)
    nexts = [
        (C_BLUE,   "Скачать реальный\nдатасет OAI"),
        (C_GREEN,  "Реализовать\nU-Net (PyTorch)"),
        (C_ORANGE, "Обучить на\nреальных масках"),
        (C_DGRAY,  "Интегрировать с\nKL-классификатором"),
    ]
    for i, (col, line) in enumerate(nexts):
        lx = Inches(2.5) + i * Inches(2.6)
        rect(s, lx, Inches(6.58), Inches(0.07), Inches(0.72), col)
        txt(s, line, lx + Inches(0.15), Inches(6.56),
            Inches(2.4), Inches(0.72), size=12, color=C_LGRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Final
# ═══════════════════════════════════════════════════════════════════════════════

def slide_final(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    rect(s, 0, 0, SW, Inches(1.2), C_DARK)
    rect(s, 0, 0, Inches(0.12), SH, C_BLUE)
    rect(s, 0, SH - Inches(0.07), SW, Inches(0.07), C_BLUE)

    txt(s, "Спасибо за внимание!",
        Inches(1.0), Inches(1.7), Inches(11.0), Inches(1.1),
        size=44, bold=True, align=PP_ALIGN.CENTER)
    rect(s, Inches(3.5), Inches(3.0), Inches(6.2), Inches(0.06), C_GREEN)
    txt(s, "Маратов Ерасыл Балканович",
        Inches(1.0), Inches(3.2), Inches(11.0), Inches(0.5),
        size=20, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
    txt(s, "7М06101  |  AIU  |  2025–2026",
        Inches(1.0), Inches(3.8), Inches(11.0), Inches(0.4),
        size=15, color=C_DGRAY, align=PP_ALIGN.CENTER)

    tags = [
        (C_RED,    "Simple Otsu"),
        (C_ORANGE, "Double-Otsu"),
        (C_BLUE,   "U-Net"),
        (C_GREEN,  "Dice Score"),
        (C_DGRAY,  "FastAPI Demo"),
    ]
    tx = Inches(0.9)
    for col, tag in tags:
        w = Inches(2.2)
        rect(s, tx, Inches(4.9), w, Inches(0.55), col)
        txt(s, tag, tx, Inches(4.95), w, Inches(0.45),
            size=13, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
        tx += w + Inches(0.15)

    txt(s, f"{TOTAL} / {TOTAL}",
        SW - Inches(1.0), SH - Inches(0.38), Inches(0.85), Inches(0.3),
        size=11, color=C_DGRAY, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# BUILD
# ═══════════════════════════════════════════════════════════════════════════════

def build():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    slide_title(prs)           # 1
    slide_problem(prs)         # 2
    slide_dataset(prs)         # 3
    slide_what_is_seg(prs)     # 4
    slide_first_approach(prs)  # 5  — Simple Otsu: brightness only
    slide_first_results(prs)   # 6  — Dice кость 0.20–0.35, хрящ = 0
    slide_why_fails(prs)       # 7  — Анализ провала
    slide_literature(prs)      # 8  — Обзор литературы
    slide_improved_method(prs) # 9  — Double-Otsu
    slide_new_results(prs)     # 10 — Новые Dice
    slide_method_choice(prs)   # 11 — Почему не нейросеть?
    slide_full_comparison(prs) # 12 — Итоговая таблица
    slide_demo(prs)            # 13 — Веб-демо
    slide_conclusions(prs)     # 14 — Выводы
    slide_final(prs)           # 15 — Спасибо

    out = Path(__file__).parent / "knee_mri_presentation_easy_explained.pptx"
    prs.save(str(out))
    print(f"Готово: {out}  ({len(prs.slides)} слайдов)")


if __name__ == "__main__":
    build()
